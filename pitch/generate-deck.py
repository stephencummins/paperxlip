"""
Paperxlip CXO Pitch Deck Generator
Generates a clean, professional PPTX for Mace Consult senior leadership.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colours
MACE_DARK = RGBColor(0x1A, 0x1A, 0x2E)
MACE_BLUE = RGBColor(0x00, 0x5C, 0x8A)
MACE_ACCENT = RGBColor(0x00, 0xA3, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MID_GREY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, colour=MACE_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 colour=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = colour
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_subtitle(slide, left, top, width, text, font_size=14, colour=MID_GREY):
    return add_text_box(slide, left, top, width, 1, text, font_size, colour)


def add_bullet_slide(slide, left, top, width, bullets, font_size=16, colour=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = colour
        p.font.name = "Calibri"
        p.space_after = spacing
    return tf


def add_accent_bar(slide, top=1.8):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(0.06), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MACE_ACCENT
    shape.line.fill.background()


def add_divider(slide, top):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Emu(12000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x44)
    shape.line.fill.background()


# ── SLIDE 1: Title ──────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text_box(slide, 0.8, 1.5, 11.7, 1.5,
             "Paperxlip", 54, MACE_ACCENT, bold=True)
add_text_box(slide, 0.8, 2.6, 11.7, 1,
             "AI Workforce Orchestration for Programme Delivery", 28, WHITE)
add_divider(slide, 3.8)
add_text_box(slide, 0.8, 4.2, 11.7, 0.5,
             "Mace Digital  |  Mace Consult  |  March 2026", 16, MID_GREY)
add_text_box(slide, 0.8, 5.0, 11.7, 0.5,
             "Confidential", 12, MID_GREY)


# ── SLIDE 2: The Problem ────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "The Problem", 36, WHITE, bold=True)
add_text_box(slide, 1.1, 1.5, 11, 0.6,
             "Institutional knowledge is fragmented, ephemeral, and walks out the door", 18, MACE_ACCENT)

bullets = [
    "Every Mace Consult project is a knowledge silo — its own SharePoint, Dataverse, Teams, email threads",
    "Critical context lives in the heads of senior people who leave or rotate",
    "No mechanism to learn from one project and apply it to the next",
    "Manual synthesis is bandwidth-limited: humans can't read everything, connect everything, remember everything",
    "When a Programme Director leaves, years of institutional knowledge disappear overnight",
    "Lessons learned decks are written, filed, and never read again",
]
add_bullet_slide(slide, 1.1, 2.4, 11, bullets, 16)


# ── SLIDE 3: The Opportunity ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "The Opportunity", 36, WHITE, bold=True)
add_text_box(slide, 1.1, 1.5, 11, 0.6,
             "The synthesis layer is the most valuable layer in enterprise software", 18, MACE_ACCENT)

bullets = [
    "Salesforce is worth $250B for owning customer data. ServiceNow is worth $200B for owning IT workflow data.",
    "The company that owns the synthesis layer across all enterprise data is worth more than both combined.",
    "OpenAI is betting $600B that they can build this. Goldman Sachs just invested in Mace Consult to build digital tools.",
    "Mace Consult sits on decades of programme delivery knowledge across hospitals, defence, transport, data centres.",
    "First mover in construction AI workforce = competitive moat that compounds with every project delivered.",
    "This isn't a tool. It's an institutional memory that gets smarter with every engagement.",
]
add_bullet_slide(slide, 1.1, 2.4, 11, bullets, 16)


# ── SLIDE 4: What is Paperxlip ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "What is Paperxlip?", 36, WHITE, bold=True)
add_text_box(slide, 1.1, 1.5, 11, 0.6,
             "An AI workforce that never forgets, never leaves, and learns from every project", 18, MACE_ACCENT)

# Left column: description
add_text_box(slide, 1.1, 2.4, 5.5, 0.5,
             "For each Mace Consult project, Paperxlip deploys a team of AI agents:", 16, WHITE)
agents = [
    "Programme Agent — coordinates all agents, produces weekly summaries, escalates to human SRO",
    "Risk Agent — monitors risk registers, flags emerging risks, finds precedent from other projects",
    "Commercial Agent — tracks variations and compensation events, drafts NEC4 responses",
    "Compliance Agent — checks RIBA stage gates, CDM, Building Safety Act, NHS requirements",
    "Knowledge Agent — continuously ingests documents, maintains searchable project knowledge",
]
add_bullet_slide(slide, 1.1, 3.0, 5.5, agents, 14, WHITE, Pt(6))

# Right column: architecture box
box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.4), Inches(5.3), Inches(4.2)
)
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

arch_lines = [
    "Mace Context Layer",
    "━━━━━━━━━━━━━━━━━━━━━━━━━━",
    "SharePoint  ·  Dataverse  ·  Teams  ·  Email",
    "",
    "         ↓  ingest + embed  ↓",
    "",
    "    Vector Store (pgvector)",
    "    Semantic search across ALL projects",
    "",
    "         ↓  query  ↓",
    "",
    "    Agent Workforce",
    "    Risk · Commercial · Programme",
    "    Compliance · Knowledge",
]
txBox = slide.shapes.add_textbox(Inches(7.5), Inches(2.6), Inches(4.9), Inches(3.8))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(arch_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.name = "Consolas"
    p.font.color.rgb = MACE_ACCENT if i in (0, 6, 7, 11, 12, 13) else RGBColor(0x99, 0x99, 0xBB)
    p.alignment = PP_ALIGN.CENTER


# ── SLIDE 5: Business Value ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "Business Value", 36, WHITE, bold=True)

# Four value boxes
values = [
    ("Knowledge\nContinuity", "Senior staff rotate or leave.\nThe knowledge stays.\nZero institutional memory loss."),
    ("Cross-Project\nPrecedent", "Risk patterns, commercial decisions,\nregulatory approaches — learned\nonce, applied everywhere."),
    ("Faster\nMobilisation", "New projects deploy an AI team\npre-loaded with Mace precedent.\nWeeks of ramp-up → hours."),
    ("Bid\nIntelligence", "Lessons from delivery feed back\ninto proposals. Win rates improve\nbecause bids are evidence-based."),
]

for i, (title, desc) in enumerate(values):
    x = 0.8 + i * 3.1
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.8), Inches(3.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)

    add_text_box(slide, x + 0.2, 2.0, 2.4, 0.9, title, 20, MACE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, 3.0, 2.4, 2.0, desc, 13, RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

# Bottom stat bar
add_divider(slide, 5.8)
stats = "Mace Consult: ~$1B revenue  ·  5,500+ specialists  ·  6 continents  ·  Decades of programme data"
add_text_box(slide, 0.8, 6.0, 11.7, 0.5, stats, 14, MID_GREY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 6: How It Works ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "How It Works", 36, WHITE, bold=True)

steps = [
    ("1. Onboard", "Connect a project's SharePoint site,\nDataverse environment, and\nTeams channels to the\nMace Context Layer."),
    ("2. Deploy", "Select a project template\n(NHP, defence, transport, data centre).\nAI agents deploy with\nrole-appropriate instructions."),
    ("3. Ingest", "Knowledge Agent continuously\nscans for new documents.\nChunks, embeds, and indexes\ninto the vector store."),
    ("4. Work", "Agents wake on schedule,\ncheck for tasks, act autonomously.\nRisk flags, commercial drafts,\ncompliance checks — all tracked."),
    ("5. Learn", "Every project enriches the\nMace Context. New projects\nlaunch with the full weight\nof institutional precedent."),
]

for i, (title, desc) in enumerate(steps):
    x = 0.5 + i * 2.5
    # Step number circle
    add_text_box(slide, x + 0.15, 1.8, 2.2, 0.5, title, 18, MACE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.15, 2.5, 2.2, 2.5, desc, 13, RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

# Arrow connectors (simple text)
for i in range(4):
    x = 2.7 + i * 2.5
    add_text_box(slide, x, 2.0, 0.5, 0.5, "→", 24, RGBColor(0x44, 0x44, 0x66), alignment=PP_ALIGN.CENTER)


# ── SLIDE 7: What We've Built ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "What We've Built", 36, WHITE, bold=True)
add_text_box(slide, 1.1, 1.5, 11, 0.6,
             "Working prototype — NHP is the first project", 18, MACE_ACCENT)

built = [
    "Open-source orchestration engine (fork of Paperclip, MIT licence) — running locally",
    "NHP company created with Programme Agent deployed on Claude Code",
    "Project template system — nhp.json defines 5 agents, org chart, goals, data sources",
    "Mace Context database schema with pgvector for semantic search (HNSW cosine index)",
    "SharePoint and Dataverse ingestor scaffolding (Microsoft Graph API)",
    "Cross-project precedent search — query any project, get relevant knowledge from all",
    "Full dashboard with real-time agent monitoring, cost tracking, and audit logs",
]
add_bullet_slide(slide, 1.1, 2.3, 11, built, 16)

add_divider(slide, 5.8)
add_text_box(slide, 0.8, 6.0, 11.7, 0.5,
             "github.com/stephencummins/paperxlip  ·  Built in 1 day  ·  Ready to demo", 14, MID_GREY, alignment=PP_ALIGN.CENTER)


# ── SLIDE 8: What We Need ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "What We Need", 36, WHITE, bold=True)

needs = [
    ("Azure AD App Registration", "Application permissions for Microsoft Graph (Sites.Read.All)\nand Dataverse API access across project tenants."),
    ("Compute & API Budget", "Cloud hosting for the orchestration engine and vector database.\nAI API costs for agent execution and embedding generation."),
    ("Pilot Project Sponsorship", "NHP as first live deployment. Access to real SharePoint data,\nrisk registers, and commercial correspondence for validation."),
    ("Mace Digital Resource", "2-3 dedicated engineers for 3 months to move from prototype\nto production-grade deployment across the NHP programme."),
]

for i, (title, desc) in enumerate(needs):
    y = 1.8 + i * 1.3
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.7), Inches(1.1)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x3A)
    box.line.color.rgb = RGBColor(0x44, 0x44, 0x66)
    add_text_box(slide, 1.1, y + 0.1, 3.5, 0.4, title, 18, MACE_ACCENT, bold=True)
    add_text_box(slide, 4.8, y + 0.15, 7.5, 0.8, desc, 14, RGBColor(0xCC, 0xCC, 0xCC))


# ── SLIDE 9: Competitive Landscape ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide, 0.8)
add_text_box(slide, 1.1, 0.7, 11, 0.8,
             "Why Now", 36, WHITE, bold=True)

bullets = [
    "Mace Consult just became independent — this is the moment to define the digital strategy",
    "Goldman Sachs backing is explicitly for 'digital tools that provide greater predictability, automation and control'",
    "OpenAI, Anthropic, and Google are racing to build the enterprise synthesis layer — but none understand construction",
    "Mace's competitors (Turner & Townsend, Arcadis, Faithful+Gould) have no AI workforce capability yet",
    "The compound advantage: every project delivered makes the system smarter — competitors can't catch up",
    "AI costs are falling 10x per year. What costs £50K/month today will cost £5K/month in 12 months.",
    "The risk is not that this doesn't work. The risk is that someone else builds it first.",
]
add_bullet_slide(slide, 1.1, 1.8, 11, bullets, 16)


# ── SLIDE 10: Close ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, 0.8, 2.0, 11.7, 1,
             "The knowledge is already there.", 36, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 3.0, 11.7, 1,
             "Paperxlip makes it usable.", 36, MACE_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_divider(slide, 4.5)

add_text_box(slide, 0.8, 5.0, 11.7, 0.5,
             "Stephen Cummins  ·  Mace Digital  ·  stephencummins@gmail.com", 16, MID_GREY, alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────
output = "/Users/stephencummins/Projects/paperxlip/pitch/Paperxlip - CXO Pitch.pptx"
prs.save(output)
print(f"Saved: {output}")
